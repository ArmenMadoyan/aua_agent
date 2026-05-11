"""Course-materials agent — generates PPTX slide decks and PDF assessments."""

from __future__ import annotations

import logging
import re
import uuid
from pathlib import Path

from fpdf import FPDF
from langchain.agents import create_agent
from langchain_core.tools import tool
from pptx import Presentation
from pptx.util import Pt

from backend.app.llm import get_creative_llm
from backend.app.prompts import COURSE_SYSTEM_PROMPT

logger = logging.getLogger(__name__)

# ── output directory ────────────────────────────────────────────────

BASE_DIR = Path(__file__).resolve().parent.parent.parent.parent
GENERATED_COURSE_DIR = BASE_DIR / "generated_course"


def ensure_generated_dir() -> None:
    GENERATED_COURSE_DIR.mkdir(parents=True, exist_ok=True)


def _safe_stem(name: str, max_len: int = 48) -> str:
    s = re.sub(r"[^\w\s-]", "", name, flags=re.UNICODE)
    s = re.sub(r"[-\s]+", "_", s.strip()).strip("_")
    return (s[:max_len] if s else "document") + f"_{uuid.uuid4().hex[:8]}"


def _artifact_url(filename: str) -> str:
    return f"/course/artifacts/{filename}"


# ── PDF builder ─────────────────────────────────────────────────────

_MODULE_DIR = Path(__file__).resolve().parent.parent
_BUNDLED_DEJAVU = _MODULE_DIR / "fonts" / "DejaVuSans.ttf"


def _pdf_font_family(pdf: FPDF) -> str:
    for font_path in (_BUNDLED_DEJAVU,):
        if not font_path.is_file():
            continue
        try:
            pdf.add_font("DejaVu", "", str(font_path))
            return "DejaVu"
        except Exception:
            continue
    try:
        import fpdf

        legacy = Path(fpdf.__file__).resolve().parent / "font" / "DejaVuSans.ttf"
        if legacy.is_file():
            pdf.add_font("DejaVu", "", str(legacy))
            return "DejaVu"
    except Exception:
        pass
    return "Helvetica"


def _ascii_fallback_text(text: str) -> str:
    out = (
        text.replace("\u2192", "->")
        .replace("\u2190", "<-")
        .replace("\u21d2", "=>")
        .replace("\u2014", "-")
        .replace("\u2013", "-")
        .replace("\u201c", '"')
        .replace("\u201d", '"')
        .replace("\u2018", "'")
        .replace("\u2019", "'")
        .replace("\u2026", "...")
    )
    return out.encode("latin-1", errors="replace").decode("latin-1")


def _build_pdf(title: str, body: str, out_path: Path, subtitle: str = "") -> None:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=18)
    family = _pdf_font_family(pdf)
    use_fallback = family == "Helvetica"
    pdf.add_page()

    def t(s: str) -> str:
        return _ascii_fallback_text(s) if use_fallback else s

    pdf.set_font(family, size=16)
    for line in title.splitlines():
        pdf.multi_cell(0, 10, t(line or " "))
    pdf.ln(4)

    if subtitle:
        pdf.set_font(family, size=12)
        pdf.multi_cell(0, 8, t(subtitle))
        pdf.ln(2)

    pdf.set_font(family, size=11)
    for para in body.split("\n\n"):
        chunk = para.strip() or " "
        pdf.multi_cell(0, 6, t(chunk))
        pdf.ln(2)

    pdf.output(str(out_path))


# ── PowerPoint builder ──────────────────────────────────────────────


def _build_powerpoint(slides: list[tuple[str, list[str]]], out_path: Path) -> None:
    prs = Presentation()
    layout = prs.slide_layouts[1]

    for title, bullets in slides:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title[:500]
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        if not bullets:
            tf.text = " "
            tf.paragraphs[0].font.size = Pt(18)
            continue
        tf.text = bullets[0][:2000]
        tf.paragraphs[0].font.size = Pt(18)
        for bullet in bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet[:2000]
            p.level = 0
            p.font.size = Pt(18)

    prs.save(str(out_path))


# ── slide text parser ───────────────────────────────────────────────


def _parse_slides_text(slides_text: str) -> list[tuple[str, list[str]]]:
    slides: list[tuple[str, list[str]]] = []
    current_title: str | None = None
    current_bullets: list[str] = []

    def flush() -> None:
        nonlocal current_title, current_bullets
        if current_title is not None:
            slides.append((current_title, current_bullets))
        current_title = None
        current_bullets = []

    for raw_line in slides_text.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        upper = line.upper()
        if upper.startswith("SLIDE:"):
            flush()
            current_title = line.split(":", 1)[1].strip() or "Slide"
            current_bullets = []
        elif line.startswith(("## ", "### ")):
            flush()
            current_title = line.lstrip("#").strip() or "Slide"
            current_bullets = []
        elif line.startswith(("- ", "* ", "\u2022 ")):
            if current_title is None:
                current_title = "Slide"
            current_bullets.append(line[2:].strip())
        elif current_title is None:
            current_title = line
        else:
            current_bullets.append(line)

    flush()

    if not slides and slides_text.strip():
        slides = [("Slides", [slides_text.strip()])]

    return slides


# ── agent class ─────────────────────────────────────────────────────


class CourseAgent:
    """Course-materials agent with PPTX and PDF generation tools."""

    def __init__(self, checkpointer=None):
        self._checkpointer = checkpointer
        self._agent = None

    @staticmethod
    @tool(
        description=(
            "Create a real PowerPoint (.pptx) file from slide content you write. "
            "Use when the user wants lecture slides or a slide deck. "
            "Argument deck_title: short name for the file. "
            "Argument slides_text: one block per slide. Start each slide with a line "
            "'SLIDE: Title Here' then bullet lines starting with '- '. "
            "Example:\n"
            "SLIDE: Introduction\n"
            "- Goal of today\n"
            "- Outline\n\n"
            "SLIDE: Next topic\n"
            "- Point one\n"
            "Returns the download path for the user."
        )
    )
    def create_powerpoint_deck(deck_title: str, slides_text: str) -> str:
        logger.info("create_powerpoint_deck called: %s", deck_title)
        ensure_generated_dir()
        slides = _parse_slides_text(slides_text)
        stem = _safe_stem(deck_title)
        filename = f"{stem}.pptx"
        out_path = GENERATED_COURSE_DIR / filename
        _build_powerpoint(slides, out_path)
        url = _artifact_url(filename)
        return (
            f"PowerPoint saved as '{filename}'. "
            f"Download from: {url} (full URL: base URL of this API + {url})"
        )

    @staticmethod
    @tool(
        description=(
            "Create a real PDF file for homework, quiz, midterm, final exam, or any written "
            "assessment or assignment. Use when the user wants homework, quizzes, exams, or "
            "problem sets as a PDF. "
            "Argument document_type: one of homework, quiz, midterm, final_exam, or other. "
            "Argument title: document title shown at the top. "
            "Argument body: full text of the assignment or exam (questions, instructions, "
            "space descriptions). Use plain text; separate major sections with a blank line. "
            "Returns the download path for the user."
        )
    )
    def create_course_pdf(document_type: str, title: str, body: str) -> str:
        logger.info("create_course_pdf called: %s / %s", document_type, title)
        ensure_generated_dir()
        stem = _safe_stem(f"{document_type}_{title}")
        filename = f"{stem}.pdf"
        out_path = GENERATED_COURSE_DIR / filename
        subtitle = f"Type: {document_type.replace('_', ' ').title()}"
        _build_pdf(title, body, out_path, subtitle=subtitle)
        url = _artifact_url(filename)
        return (
            f"PDF saved as '{filename}'. "
            f"Download from: {url} (full URL: base URL of this API + {url})"
        )

    def build(self):
        logger.info("Building course agent")
        self._agent = create_agent(
            model=get_creative_llm(),
            tools=[self.create_powerpoint_deck, self.create_course_pdf],
            system_prompt=COURSE_SYSTEM_PROMPT,
            checkpointer=self._checkpointer,
            name="course_agent",
        )
        return self._agent

    def get_agent(self):
        if self._agent is None:
            self.build()
        return self._agent
