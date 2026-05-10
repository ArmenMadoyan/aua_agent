from pathlib import Path

from fpdf import FPDF
from pptx import Presentation
from pptx.util import Pt

_MODULE_DIR = Path(__file__).resolve().parent
_BUNDLED_DEJAVU = _MODULE_DIR / "fonts" / "DejaVuSans.ttf"


def _pdf_font_family(pdf: FPDF) -> str:
    """Prefer bundled DejaVu (full Unicode); fpdf2 no longer ships Helvetica-safe Unicode."""
    for font_path in (_BUNDLED_DEJAVU,):
        if not font_path.is_file():
            continue
        try:
            pdf.add_font("DejaVu", "", str(font_path))
            return "DejaVu"
        except Exception:
            continue
    try:
        import fpdf

        legacy = Path(fpdf.__file__).resolve().parent / "font" / "DejaVuSans.ttf"
        if legacy.is_file():
            pdf.add_font("DejaVu", "", str(legacy))
            return "DejaVu"
    except Exception:
        pass
    return "Helvetica"


def _ascii_fallback_text(text: str) -> str:
    """Last resort if no Unicode font (avoids fpdf crash on arrows, smart quotes, etc.)."""
    out = (
        text.replace("→", "->")
        .replace("←", "<-")
        .replace("⇒", "=>")
        .replace("—", "-")
        .replace("–", "-")
        .replace(""", '"')
        .replace(""", '"')
        .replace("'", "'")
        .replace("'", "'")
        .replace("…", "...")
    )
    return out.encode("latin-1", errors="replace").decode("latin-1")


def build_powerpoint(slides: list[tuple[str, list[str]]], out_path: Path) -> None:
    prs = Presentation()
    layout = prs.slide_layouts[1]

    for title, bullets in slides:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title[:500]
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        if not bullets:
            tf.text = " "
            tf.paragraphs[0].font.size = Pt(18)
            continue
        tf.text = bullets[0][:2000]
        tf.paragraphs[0].font.size = Pt(18)
        for bullet in bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet[:2000]
            p.level = 0
            p.font.size = Pt(18)

    prs.save(str(out_path))


def build_pdf(title: str, body: str, out_path: Path, subtitle: str = "") -> None:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=18)
    family = _pdf_font_family(pdf)
    use_fallback = family == "Helvetica"
    pdf.add_page()

    def t(s: str) -> str:
        return _ascii_fallback_text(s) if use_fallback else s

    pdf.set_font(family, size=16)
    for line in title.splitlines():
        pdf.multi_cell(0, 10, t(line or " "))
    pdf.ln(4)

    if subtitle:
        pdf.set_font(family, size=12)
        pdf.multi_cell(0, 8, t(subtitle))
        pdf.ln(2)

    pdf.set_font(family, size=11)
    for para in body.split("\n\n"):
        chunk = para.strip() or " "
        pdf.multi_cell(0, 6, t(chunk))
        pdf.ln(2)

    pdf.output(str(out_path))
